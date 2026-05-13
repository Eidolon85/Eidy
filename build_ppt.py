"""Generate an editable PPTX that recreates the 5 source slides verbatim.

Layout/styling approximates the originals; all text is editable.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- helpers ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_text(slide, x, y, w, h, text, *, size=11, bold=False, color=RGBColor(0x20, 0x20, 0x20),
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None, font="微软雅黑"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.shadow.inherit = False
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line_text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box


def add_rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             fill=None, line=None):
    """runs is a list of paragraphs; each paragraph is a list of (text, opts) tuples.
    opts may include size, bold, color, italic, font, align."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.shadow.inherit = False
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for text, opts in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(opts.get("size", 11))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.name = opts.get("font", "微软雅黑")
            color = opts.get("color", RGBColor(0x20, 0x20, 0x20))
            r.font.color.rgb = color
    return box


# Common colors
DARK = RGBColor(0x1F, 0x2A, 0x44)
NAVY = RGBColor(0x0B, 0x2A, 0x6B)
BLUE = RGBColor(0x21, 0x6B, 0xD0)
LIGHT_BLUE = RGBColor(0x3E, 0xA7, 0xE6)
TEAL = RGBColor(0x1F, 0x8F, 0x8A)
GREEN = RGBColor(0x2E, 0xA5, 0x6A)
PURPLE = RGBColor(0x6A, 0x3D, 0xA8)
ORANGE = RGBColor(0xE2, 0x88, 0x18)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK = RGBColor(0x0E, 0x16, 0x36)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============================================================
# Slide 1: AI Case & Offering 总结
# ============================================================
s1 = prs.slides.add_slide(blank)

add_text(s1, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7),
         "AI Case & Offering 总结", size=32, bold=True, color=DARK)
add_text(s1, Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.4),
         "我们在过去一年在多个业务和技术维度成功落地了多个典型AI项目",
         size=14, color=GREY)
# horizontal rule
add_text(s1, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.03),
         "", fill=RGBColor(0xCF, 0xD8, 0xE3))

# Left vertical labels
add_text(s1, Inches(0.4), Inches(1.55), Inches(0.45), Inches(4.2),
         "AI Enable Business", size=10, bold=True, color=DARK,
         fill=RGBColor(0xE6, 0xE0, 0xEC), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Column headers
add_text(s1, Inches(1.0),  Inches(1.55), Inches(3.9), Inches(0.4),
         "Content Generation", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s1, Inches(5.0),  Inches(1.55), Inches(3.9), Inches(0.4),
         "Efficiency Improvement", size=14, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1, Inches(9.0),  Inches(1.55), Inches(3.9), Inches(0.4),
         "WoW Innovation", size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
# underline bars
add_text(s1, Inches(1.0), Inches(1.95), Inches(3.9), Inches(0.04), "", fill=NAVY)
add_text(s1, Inches(5.0), Inches(1.95), Inches(3.9), Inches(0.04), "", fill=LIGHT_BLUE)
add_text(s1, Inches(9.0), Inches(1.95), Inches(3.9), Inches(0.04), "", fill=TEAL)

# Column 1: Content Generation - single big block
add_rich(s1, Inches(1.0), Inches(2.05), Inches(3.9), Inches(3.7), [
    [("\n\nAI 效果图助手", {"size": 18, "bold": True, "color": DARK})],
    [("全员可用的生图工具与智能图像处理", {"size": 11, "color": GREY})],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=NAVY)

# Column 2: Efficiency Improvement
eff_items = [
    ("AI 目标识别", "实现高性能 OCR 和图章检测，满足预定义基准", "J&J"),
    ("AI 文献搜索：百济", "人机互动全链路赋能销售代表", "BeiGene"),
    ("AI日志分析", "轮询实现日志的自动化推理检测与异常日志实时告警", "Porsche"),
    ("AI 合同管理：沃尔沃/越秀", "“智能秒级”响应，自动解析合同条款", "YX"),
    ("智能客服系统", "7*24快速响应客户需求并自动记录到业务系统", "Porsche"),
]
y = 2.05
for title, desc, tag in eff_items:
    add_rich(s1, Inches(5.0), Inches(y), Inches(3.4), Inches(0.7), [
        [(title, {"size": 13, "bold": True, "color": DARK})],
        [(desc, {"size": 9, "color": GREY})],
    ], align=PP_ALIGN.CENTER)
    add_text(s1, Inches(8.45), Inches(y+0.1), Inches(0.45), Inches(0.45),
             tag, size=8, color=WHITE, fill=RGBColor(0xC8, 0x10, 0x2E),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.74

# Column 3: WoW Innovation
wow_items = [
    ("Deep Research 混合检索引擎", "集成知识库 + 外部网站搜索 + 内部资产 Confluence 为一体，整合多源检索与多轮优化能力", "BMW"),
    ("AI智能陪练", "线上实时陪练，数据驱动优化效果", ""),
    ("AI 知识库", "打造企业“智能知识中枢”，实现24h智能服务", ""),
    ("Machine Learning Ops", "实现可复现、可追溯、可部署的机器学习工程化流程", "J&J"),
]
y = 2.05
heights = [1.0, 0.8, 0.8, 1.0]
for (title, desc, tag), hh in zip(wow_items, heights):
    add_rich(s1, Inches(9.0), Inches(y), Inches(3.4), Inches(hh), [
        [(title, {"size": 13, "bold": True, "color": DARK})],
        [(desc, {"size": 9, "color": GREY})],
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if tag:
        add_text(s1, Inches(12.4), Inches(y+0.1), Inches(0.5), Inches(0.4),
                 tag, size=8, color=WHITE, fill=RGBColor(0x1A, 0x4D, 0x9E),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += hh

# AI Evolve Workspace left label
add_text(s1, Inches(0.4), Inches(5.85), Inches(0.45), Inches(1.5),
         "AI Evolve Workspace", size=9, bold=True, color=DARK,
         fill=RGBColor(0xE6, 0xE0, 0xEC), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Product X strip
add_rich(s1, Inches(1.0), Inches(5.85), Inches(11.9), Inches(0.75), [
    [("Product X", {"size": 14, "bold": True, "color": DARK})],
    [("以大模型等AI技术为驱动，以提高软件研发运营智能化水平为导向，以提质增效为目标的新一代智能化软件工程",
      {"size": 9, "color": GREY})],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=ORANGE)

# Product X sub-modules
mods = ["智能项目管理", "智能需求设计", "智能开发", "智能测试", "部署和运维"]
mw = 11.9 / 5
for i, m in enumerate(mods):
    add_text(s1, Inches(1.0 + i*mw + 0.05), Inches(6.65), Inches(mw-0.1), Inches(0.32),
             m, size=10, color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             fill=LIGHT_GREY)

# 智擎 strip
add_rich(s1, Inches(1.0), Inches(7.05), Inches(11.9), Inches(0.4), [
    [("智擎", {"size": 12, "bold": True, "color": DARK}),
     ("    面向企业客户的智能体应用开发平台，以企业知识库为内核，帮助用户个性化定制基于大模型的智能体应用，并快速集成和部署到不同的平台",
      {"size": 8, "color": GREY})],
], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line=ORANGE)

# Footer
add_text(s1, Inches(10.5), Inches(7.25), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 2: AI ERP 合同赋能
# ============================================================
s2 = prs.slides.add_slide(blank)

# Background split: left light, right dark
add_text(s2, Inches(0), Inches(0), Inches(5.0), SLIDE_H, "",
         fill=RGBColor(0xF4, 0xF6, 0xFA))
add_text(s2, Inches(5.0), Inches(0), Inches(8.333), SLIDE_H, "",
         fill=BG_DARK)

# Left header
add_text(s2, Inches(0.4), Inches(0.3), Inches(4.4), Inches(0.7),
         "AI ERP 合同赋能", size=28, bold=True, color=DARK)
add_text(s2, Inches(0.4), Inches(1.05), Inches(4.4), Inches(0.9),
         "让风险洞察“先知先觉”，“智能秒级”响应，自动解析合同条款",
         size=14, bold=True, color=DARK)

# Left bullets
left_sections = [
    ("处理速度大幅提升", [
        "传统人工审核：每份合同需要2-4小时",
        "AI智能检测：单份合同检测≤30秒",
        "效率提升：百倍速度提升",
        "时间节约：每个工作日可节省大量人工时间",
    ]),
    ("工作流程自动化", [
        "6大智能功能模块：租户资质、合同条款、价格分析、免租期、信用评估、舆情监控",
        "并行数据查询：同步调用多个外部数据源",
        "智能报告生成：自动生成分析报告",
        "减少人工干预：90%以上流程实现自动化",
    ]),
    ("风险管控价值", [
        "风险识别能力提升：多维度风险评估：企业信用、经营异常、限制高消费、舆情风险",
        "风险预警机制：提前识别潜在风险租户",
        "预计避免损失：每年可避免百万级别风险损失",
        "合规性提升：政策对标检查：自动对比租赁政策和历史合同",
    ]),
]
y = 2.05
for title, items in left_sections:
    add_text(s2, Inches(0.4), Inches(y), Inches(4.4), Inches(0.32),
             title, size=12, bold=True, color=BLUE)
    y += 0.35
    bullet_text = "\n".join("• " + it for it in items)
    h = 0.22 * len(items) + 0.05
    add_text(s2, Inches(0.5), Inches(y), Inches(4.3), Inches(h + 0.4),
             bullet_text, size=8.5, color=DARK)
    y += h + 0.15

# Right header
add_text(s2, Inches(5.3), Inches(0.3), Inches(4.0), Inches(0.6),
         "解决方案", size=22, bold=True, color=LIGHT_BLUE)

# Right hex-icon sections
def hex_icon(slide, x, y, label, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, Inches(0.7), Inches(0.7))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame; tf.margin_left=tf.margin_right=Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "微软雅黑"

right_sections = [
    ("场景", "JDE系统发起合同检测请求，传入合同基本信息和租户信息\n根据租户类型和分析需求，调用启信宝、信用查询、舆情监控等外部API"),
    ("效果", "通过AI大模型执行6大功能模块的智能分析，帮助企业快速决策，提高检测时间效率，识别潜在风险。"),
    ("亮点", "AI：DIFY平台通过通义大模型执行6大功能模块的智能分析\n多系统集成：启信宝、信用查询、舆情监控等外部API，汇总多维度数据，作为AI上下文输入"),
]
y = 1.1
for label, body in right_sections:
    hex_icon(s2, Inches(5.3), Inches(y), label, BLUE)
    add_text(s2, Inches(6.1), Inches(y), Inches(3.6), Inches(0.4),
             label, size=14, bold=True, color=LIGHT_BLUE)
    add_text(s2, Inches(6.1), Inches(y+0.4), Inches(3.6), Inches(1.4),
             body, size=10, color=WHITE)
    y += 1.85

# Right placeholder for the analysis report visual (text representation)
add_text(s2, Inches(10.0), Inches(1.1), Inches(3.2), Inches(5.3),
         "[ 合同分析报告示例 ]\n\n— 多维度数据汇总表\n— 风险评分明细\n— 合规检查项",
         size=9, color=RGBColor(0xFF, 0xB3, 0x4D), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         fill=RGBColor(0x14, 0x1F, 0x4A), line=ORANGE)


# ============================================================
# Slide 3: AI Product X
# ============================================================
s3 = prs.slides.add_slide(blank)

# Left side header
add_text(s3, Inches(0.4), Inches(0.4), Inches(4.0), Inches(0.7),
         "AI Product X", size=30, bold=True, color=DARK)
add_text(s3, Inches(0.4), Inches(1.1), Inches(4.0), Inches(0.8),
         "为软件全流程赋能\n实现降本增效与加速创新",
         size=18, bold=True, color=RGBColor(0xB8, 0xA9, 0x6E))
add_text(s3, Inches(0.4), Inches(2.4), Inches(4.0), Inches(2.6),
         "AI PRODUCT X智能辅助平台，利用AI技术重塑软件开发生命周期(SDLC)的每一个环节，"
         "通过需求智能化、开发智能化、测试智能化、运维智能化保障交付付质量，加速业务价值实现。",
         size=11, color=DARK)

# Right header
add_text(s3, Inches(4.6), Inches(0.4), Inches(8.3), Inches(0.5),
         "解决方案", size=20, bold=True, color=LIGHT_BLUE)
add_text(s3, Inches(4.6), Inches(0.95), Inches(8.3), Inches(0.6),
         "AI4SE(AI for Software Engineering)：是以大模型等AI技术为驱动，"
         "以提高软件研发运营智能化水平为导向，以提质增效为目标的新一代智能化软件工程。",
         size=10, color=DARK)

# Matrix
mat_x = Inches(4.6); mat_y = Inches(1.7)
col_w = Inches(1.55); row_h = Inches(0.32)
headers_top = ["", "智能化水平", "应用成熟度", "", "应用成效"]
# Top group header bar
for i, h in enumerate(headers_top):
    if not h: continue
    add_text(s3, mat_x + col_w*(i), mat_y, col_w, Inches(0.35),
             h, size=10, bold=True, color=WHITE, fill=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Column headers
col_titles = ["应用能力", "智能项目管理", "智能需求设计", "智能开发", "智能测试", "部署和运维"]
col_colors = [DARK, LIGHT_BLUE, GREEN, GREEN, GREEN, LIGHT_BLUE]
cw = Inches(1.4)
for i, (t, c) in enumerate(zip(col_titles, col_colors)):
    add_text(s3, mat_x + cw*i, mat_y + Inches(0.4), cw - Inches(0.05), Inches(0.35),
             t, size=10, bold=True, color=WHITE if i>0 else DARK,
             fill=(c if i>0 else None),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Matrix rows
rows = [
    ["平台\n工具\n能力", "知识库能力", "需求分析能力", "代码生成能力", "测试分析能力", "配置生成能力"],
    ["",                "流程调度能力", "需求生成能力", "代码补全能力", "测试用例生成能力", "代码检测能力"],
    ["",                "质量管理能力", "架构设计能力", "单测能力",     "测试脚本生成能力", "智能监控能力"],
    ["",                "效能分析能力", "UI/接口设计能力","代码注释能力","智能测试能力",   "SQL脚本生成能力"],
    ["",                "过程检查能力", "评审检查能力", "编码辅助能力", "缺陷分析能力",     "故障分析/纠正能力"],
]
row_colors = [PURPLE, BLUE, BLUE, BLUE, BLUE, BLUE]
y0 = mat_y + Inches(0.78)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        if ci == 0:
            if ri == 0:
                add_text(s3, mat_x, y0, cw - Inches(0.05), Inches(0.34*5),
                         "平台\n工具\n能力", size=10, bold=True, color=DARK,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            continue
        color = PURPLE if (ri == 0 and ci == 1) else BLUE
        add_text(s3, mat_x + cw*ci, y0 + Inches(0.34*ri), cw - Inches(0.05), Inches(0.30),
                 val, size=9, color=WHITE, fill=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom bullets
bullets = [
    ("Agent首页功能完成：", "实现Agent模块首页，提供关键数据概览与核心功能入口，提升操作便捷性。"),
    ("Subtask Agent对话历史快速操作与导航：", "为Subtask Agent同步实现对话历史的快速操作与导航功能，保持体验一致性。"),
    ("技术方案Agent对话历史快速操作与导航：", "为技术方案Agent增加对话历史的快速操作（如删除、分享）及导航，优化用户交互效率。"),
    ("历史记录跳转功能：", "支持从历史记录列表直接跳转到对应对话详情，减少操作步骤，提升导航效率。"),
]
by = 5.9
positions = [(4.6, by), (9.0, by), (4.6, by+0.85), (9.0, by+0.85)]
order = [0, 2, 1, 3]
for slot, idx in zip(positions, order):
    bx, by_ = slot
    title, body = bullets[idx]
    add_rich(s3, Inches(bx), Inches(by_), Inches(4.2), Inches(0.8), [
        [("• ", {"size": 10, "color": LIGHT_BLUE, "bold": True}),
         (title, {"size": 10, "bold": True, "color": DARK}),
         (body, {"size": 9, "color": DARK})],
    ])


# ============================================================
# Slide 4: 我们的AI赋能开发蓝图
# ============================================================
s4 = prs.slides.add_slide(blank)

add_text(s4, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6),
         "我们的AI赋能开发蓝图", size=28, bold=True, color=DARK)
add_text(s4, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
         "从协同到自主——研发效能成熟度四级跃迁模型",
         size=13, color=GREY)

# AI x 效能实践 area - light blue box
add_text(s4, Inches(0.4), Inches(1.4), Inches(10.4), Inches(3.0), "",
         fill=RGBColor(0xE8, 0xF1, 0xFB))
# Left vertical label
add_text(s4, Inches(0.45), Inches(1.5), Inches(0.5), Inches(2.8),
         "AI\n×\n效能\n实践", size=10, bold=True, color=DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Stage badges row
def circle_badge(slide, x, y, text, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.45), Inches(0.45))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Arial"

# top stages
stages = [
    ("L0", "人工开发", DARK),
    ("L1", "AI辅助（Copilot）", NAVY),
    ("L2", "AI协同（Agent）", RGBColor(0xE5, 0x6B, 0x6B)),
    ("L3", "AI自主（Agentic）", ORANGE),
]
xs = [1.2, 3.0, 5.5, 8.5]
add_text(s4, Inches(1.0), Inches(1.5), Inches(1.4), Inches(0.35),
         "四阶级实践演进", size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
for x, (lbl, name, col) in zip(xs, stages):
    circle_badge(s4, Inches(x), Inches(1.45), lbl, col)
    add_text(s4, Inches(x+0.5), Inches(1.5), Inches(2.0), Inches(0.35),
             name, size=10, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# Org-level row
add_text(s4, Inches(1.0), Inches(2.0), Inches(1.5), Inches(0.4),
         "组织级 效能提升", size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(2.5), Inches(2.05), Inches(0.7), Inches(0.35),
         "精益提效", size=9, color=DARK, fill=RGBColor(0xCF,0xE0,0xF4),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(3.3), Inches(2.05), Inches(5.6), Inches(0.35),
         "AI First/AI Ready研发组织转型", size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(9.0), Inches(2.05), Inches(1.7), Inches(0.35),
         "AI Native研发组织转型", size=9, bold=True, color=DARK, fill=RGBColor(0xCF,0xE0,0xF4),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Team-level row
add_text(s4, Inches(1.0), Inches(2.5), Inches(1.5), Inches(0.4),
         "团队级 研发模式", size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(2.5), Inches(2.55), Inches(0.7), Inches(0.35),
         "敏捷开发", size=9, color=DARK, fill=RGBColor(0xCF,0xE0,0xF4),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(3.3), Inches(2.55), Inches(2.7), Inches(0.35),
         "AI增强型 敏捷开发", size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(6.1), Inches(2.55), Inches(2.8), Inches(0.35),
         "AI驱动型 敏捷开发", size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(9.0), Inches(2.55), Inches(1.7), Inches(0.35),
         "AI Native开发", size=9, bold=True, color=DARK, fill=RGBColor(0xCF,0xE0,0xF4),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Personal-level row
add_text(s4, Inches(1.0), Inches(3.0), Inches(1.5), Inches(0.4),
         "个人级 开发方法", size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(2.5), Inches(3.05), Inches(0.7), Inches(0.35),
         "标准开发", size=9, color=DARK, fill=RGBColor(0xCF,0xE0,0xF4),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(3.3), Inches(3.05), Inches(2.7), Inches(0.35),
         "AI辅助编码", size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(6.1), Inches(3.05), Inches(2.8), Inches(0.35),
         "AI辅助/协同开发", size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(9.0), Inches(3.05), Inches(1.7), Inches(0.35),
         "AI自主开发", size=9, bold=True, color=DARK, fill=RGBColor(0xCF,0xE0,0xF4),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Right side: 报告 / 模型 / 指标
add_text(s4, Inches(10.95), Inches(1.45), Inches(0.5), Inches(0.4),
         "报告", size=9, color=DARK)
add_text(s4, Inches(11.5), Inches(1.45), Inches(1.6), Inches(0.4),
         "效能洞察报告", size=9, color=DARK, fill=LIGHT_GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(10.95), Inches(1.95), Inches(0.5), Inches(0.4),
         "模型", size=9, color=DARK)
add_text(s4, Inches(11.5), Inches(1.95), Inches(1.6), Inches(0.4),
         "研发效能 模型\nAI研发成熟度 模型", size=8, color=DARK, fill=LIGHT_GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(10.95), Inches(2.5), Inches(0.5), Inches(0.4),
         "指标", size=9, color=DARK)
add_text(s4, Inches(11.5), Inches(2.5), Inches(0.8), Inches(1.4),
         "需求交付\n开发效率\n工程效率\n……", size=8, color=DARK, fill=LIGHT_GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(12.3), Inches(2.5), Inches(0.8), Inches(1.4),
         "组织级\n业务线级\n团队级\n个人级\n……", size=8, color=DARK, fill=LIGHT_GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# AI x 效能平台 area (bottom)
add_text(s4, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.85), "",
         fill=RGBColor(0xF6, 0xF8, 0xFB))
add_text(s4, Inches(0.45), Inches(4.6), Inches(0.5), Inches(2.6),
         "AI\n×\n效能\n平台", size=10, bold=True, color=DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ProductX label
add_text(s4, Inches(1.0), Inches(4.55), Inches(1.5), Inches(0.6),
         "ProductX\n智能效率平台", size=10, bold=True, color=BLUE,
         anchor=MSO_ANCHOR.MIDDLE)

# stage chevron row
phases = ["需求阶段", "开发阶段", "测试阶段", "集成&发布阶段", "运维阶段", "度量与改进阶段"]
phase_widths = [1.4, 2.6, 1.4, 1.7, 1.4, 1.9]
px = 2.6
for ph, pw in zip(phases, phase_widths):
    sh = s4.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(px), Inches(4.6), Inches(pw), Inches(0.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
    tf = sh.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ph; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "微软雅黑"
    px += pw - 0.1

# AI能力 row
add_text(s4, Inches(1.0), Inches(5.15), Inches(1.5), Inches(0.4),
         "AI能力", size=10, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(2.6), Inches(5.15), Inches(7.4), Inches(0.35),
         "Native交付流（Flow）", size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(10.1), Inches(5.15), Inches(1.5), Inches(0.35),
         "On Call/运维 Agent", size=9, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(11.7), Inches(5.15), Inches(1.2), Inches(0.35),
         "度量Agent", size=9, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s4, Inches(2.6), Inches(5.55), Inches(4.0), Inches(0.35),
         "需求&技术设计Agent   Code Agent   CR Agent   UT Agent   ……",
         size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(6.7), Inches(5.55), Inches(1.5), Inches(0.35),
         "测试Agent", size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(8.3), Inches(5.55), Inches(1.7), Inches(0.35),
         "研发流程Agent", size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# 研发工具链 row
add_text(s4, Inches(1.0), Inches(6.0), Inches(1.5), Inches(0.8),
         "研发\n工具链", size=10, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(2.6), Inches(6.0), Inches(1.4), Inches(0.7),
         "协作平台\nProductX", size=10, bold=True, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(4.0), Inches(6.0), Inches(6.0), Inches(0.35),
         "多端交付平台：ProductX-Web/Lingma-Desktop", size=10, bold=True,
         color=WHITE, fill=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(4.0), Inches(6.35), Inches(6.0), Inches(0.35),
         "质量平台：ProductX-Test", size=10, bold=True,
         color=WHITE, fill=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(10.1), Inches(6.0), Inches(1.5), Inches(0.7),
         "运维平台\nProductX-Ops", size=9, bold=True, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(11.7), Inches(6.0), Inches(1.2), Inches(0.7),
         "度量平台", size=9, bold=True, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# AI底座 row
add_text(s4, Inches(1.0), Inches(6.85), Inches(1.5), Inches(0.4),
         "AI底座", size=10, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, Inches(2.6), Inches(6.85), Inches(10.3), Inches(0.4),
         "Qwen模型（研发数据、Agent记忆与自学系统）",
         size=10, bold=True, color=WHITE, fill=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# Slide 5: AI 技术能力矩阵
# ============================================================
s5 = prs.slides.add_slide(blank)

add_text(s5, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55),
         "AI 技术能力矩阵", size=26, bold=True, color=DARK)

# Two top section headers
add_text(s5, Inches(0.4), Inches(0.85), Inches(5.8), Inches(0.4),
         "🧠  大语言模型层", size=12, bold=True, color=WHITE, fill=BLUE,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s5, Inches(6.3), Inches(0.85), Inches(6.6), Inches(0.4),
         "⚙  Agent 开发与编排", size=12, bold=True, color=WHITE, fill=PURPLE,
         anchor=MSO_ANCHOR.MIDDLE)

# LLM layer grid (4 rows x 3 cols)
llm_items = [
    ["OpenAI", "deepseek", "Ollama"],
    ["通义千问", "文心一言", "Gemini"],
    ["Azure OpenAI Service", "aws Bedrock", "智谱·AI"],
    ["Hugging Face", "OpenLLM", "01.AI"],
]
gx, gy = 0.4, 1.3
cw_, rh_ = 1.9, 0.55
for r, row in enumerate(llm_items):
    for c, item in enumerate(row):
        add_text(s5, Inches(gx + c*cw_), Inches(gy + r*rh_), Inches(cw_ - 0.05), Inches(rh_ - 0.05),
                 item, size=11, bold=True, color=DARK, fill=LIGHT_GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Agent dev cols
agent_x = 6.3
agent_cw = 2.2
# headers
add_text(s5, Inches(agent_x), Inches(1.3), Inches(agent_cw-0.05), Inches(0.3),
         "框架", size=10, bold=True, color=WHITE, fill=RGBColor(0xE8,0x6B,0x9C),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s5, Inches(agent_x+agent_cw), Inches(1.3), Inches(agent_cw-0.05), Inches(0.3),
         "平台", size=10, bold=True, color=WHITE, fill=RGBColor(0xE8,0x6B,0x9C),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s5, Inches(agent_x+agent_cw*2), Inches(1.3), Inches(agent_cw-0.05), Inches(0.3),
         "协议/观测", size=10, bold=True, color=WHITE, fill=RGBColor(0xE8,0x6B,0x9C),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

framework_items = [
    ("LangChain / LangGraph", "链式·记忆·工具调用"),
    ("AutoGen (微软)", "多Agent对话协作"),
    ("CrewAI", "角色Agent·任务分派"),
    ("Semantic Kernel", "企业级·插件架构"),
]
platform_items = [
    ("智擎（凯捷）/Dify / Coze", "可视化编排·企业Bot"),
    ("钉钉 AI / 飞书 AI", "企微集成·办公自动化"),
    ("阿里云百炼", "DashScope·ModelScope"),
    ("百度智能体", "文心·知识库集成"),
]
protocol_items = [
    ("MCP (Model Context)", "Anthropic·工具互操作"),
    ("Function Calling / A2A", "OpenAI·Google标准接口"),
    ("LangSmith / Helicone", "调用追踪·成本分析"),
]
def draw_items(slide, x, y, w, items):
    yy = y
    for name, desc in items:
        add_rich(slide, Inches(x), Inches(yy), Inches(w), Inches(0.5), [
            [("▸ ", {"size": 10, "color": PURPLE, "bold": True}),
             (name, {"size": 10, "bold": True, "color": DARK})],
            [(desc, {"size": 8, "color": GREY})],
        ])
        yy += 0.55

draw_items(s5, agent_x, 1.65, agent_cw-0.05, framework_items)
draw_items(s5, agent_x+agent_cw, 1.65, agent_cw-0.05, platform_items)
draw_items(s5, agent_x+agent_cw*2, 1.65, agent_cw-0.05, protocol_items)

# Middle section: AI 开发工具 & MLOps
add_text(s5, Inches(0.4), Inches(3.95), Inches(6.8), Inches(0.35),
         "💻  AI 开发工具 & 编码助手", size=12, bold=True, color=WHITE,
         fill=RGBColor(0x2C, 0x7A, 0x4F), anchor=MSO_ANCHOR.MIDDLE)
add_text(s5, Inches(7.3), Inches(3.95), Inches(5.6), Inches(0.35),
         "📊  MLOps · 评测 · 观测", size=12, bold=True, color=WHITE,
         fill=RGBColor(0x2C, 0x7A, 0x4F), anchor=MSO_ANCHOR.MIDDLE)

dev_tools = [
    ("GitHub Copilot", "代码补全·单元测试"),
    ("Cursor / windsurf", "AI IDE·整文件重构"),
    ("Claude (CLI)", "全栈·架构设计"),
    ("Codex", "来自OpenAI·goal命令"),
    ("Lingma", "国内合规·阿里系"),
]
mlops_tools = [
    ("vLLM / Ollama", "高吞吐推理·本地部署"),
    ("Axolotl / TRL", "分布式训练·RLHF/DPO"),
    ("LM Arena / LMSYS", "开源模型评测基准"),
    ("OpenCompass", "国产评测·中文能力"),
    ("FlagEval", "多维度·国产大模型评测"),
]
dt_w = 6.8/5
for i, (n, d) in enumerate(dev_tools):
    add_rich(s5, Inches(0.4 + i*dt_w), Inches(4.35), Inches(dt_w-0.05), Inches(0.85), [
        [(n, {"size": 10, "bold": True, "color": DARK})],
        [(d, {"size": 8, "color": GREY})],
    ], fill=RGBColor(0xEC, 0xF6, 0xEF))
mt_w = 5.6/5
for i, (n, d) in enumerate(mlops_tools):
    add_rich(s5, Inches(7.3 + i*mt_w), Inches(4.35), Inches(mt_w-0.05), Inches(0.85), [
        [(n, {"size": 10, "bold": True, "color": DARK})],
        [(d, {"size": 8, "color": GREY})],
    ], fill=RGBColor(0xEC, 0xF6, 0xEF))

# Bottom section: AI 工程化实践
add_text(s5, Inches(0.4), Inches(5.35), Inches(12.5), Inches(0.35),
         "🛠  AI 工程化实践", size=12, bold=True, color=WHITE, fill=PURPLE,
         anchor=MSO_ANCHOR.MIDDLE)

practices = [
    ("SDD · 规格驱动开发",
     "形式化规格 → AI 精准生成代码\n降低返工，0规格偏差交付",
     "FSD → Code · 双向追溯 · AI Review",
     RGBColor(0xF2, 0xA7, 0x3C)),
    ("TDD · 测试驱动 AI",
     "AI 自动生成测试用例\n覆盖边界条件，回归自动化",
     "用例生成 · Mock 数据 · 变异测试",
     RGBColor(0xE2, 0x6B, 0x4F)),
    ("MLOps · 模型运维",
     "训练→部署→监控→迭代\n全生命周期规模化管理",
     "CI/CD for ML · A/B Testing · 漂移检测",
     RGBColor(0x4F, 0x9F, 0x55)),
    ("Harness · 评测驱动",
     "MMLU / HumanEval / MMLU-Pro\n国产: OpenCompass / FlagEval",
     "能力评测 · 红队对抗 · 安全对齐",
     RGBColor(0xE0, 0x9F, 0x4F)),
    ("RAG · 检索增强",
     "向量数据库 + 大模型\n企业知识库问答事实增强",
     "Embedding · 向量检索 · 知识库",
     RGBColor(0x3A, 0x8F, 0xC8)),
    ("Agent · 自主编排",
     "多步推理·工具调用\nReAct / Plan-and-Execute",
     "ReAct · Tool Use · 记忆管理",
     RGBColor(0x8E, 0x5C, 0xC2)),
]
pw = 12.5/6
for i, (title, body, foot, col) in enumerate(practices):
    bx = 0.4 + i*pw
    add_text(s5, Inches(bx), Inches(5.75), Inches(pw-0.05), Inches(0.32),
             title, size=10, bold=True, color=WHITE, fill=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s5, Inches(bx), Inches(6.07), Inches(pw-0.05), Inches(0.75),
             body, size=8, color=DARK, fill=RGBColor(0xFA,0xFA,0xFA))
    add_text(s5, Inches(bx), Inches(6.85), Inches(pw-0.05), Inches(0.45),
             foot, size=8, color=WHITE, fill=col,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_text(s5, Inches(10.5), Inches(7.30), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 6: AI 陪练
# ============================================================
s6 = prs.slides.add_slide(blank)

# Left light panel, right dark panel
add_text(s6, Inches(0), Inches(0), Inches(4.5), SLIDE_H, "",
         fill=RGBColor(0xF4, 0xF6, 0xFA))
add_text(s6, Inches(4.5), Inches(0), Inches(8.833), SLIDE_H, "",
         fill=BG_DARK)

# Left title
add_text(s6, Inches(0.4), Inches(0.3), Inches(4.0), Inches(0.7),
         "AI 陪练", size=32, bold=True, color=DARK)
add_text(s6, Inches(0.4), Inches(1.1), Inches(4.0), Inches(1.0),
         "线上实时陪练\n+\n数据驱动优化",
         size=16, color=ORANGE, align=PP_ALIGN.CENTER)

# Left business value
add_text(s6, Inches(0.4), Inches(2.7), Inches(4.0), Inches(0.4),
         "业务价值", size=16, bold=True, color=ORANGE)

bv6 = [
    "提升员工培训效率，随时随地通过移动端参与培训",
    "显著降低长期培训成本",
    "员工的能力数据透明化，培训提升更具针对性",
]
y = 3.2
for item in bv6:
    add_text(s6, Inches(0.4), Inches(y), Inches(0.1), Inches(0.5),
             "", fill=LIGHT_BLUE)
    add_text(s6, Inches(0.55), Inches(y), Inches(3.8), Inches(0.8),
             item, size=11, color=DARK)
    y += 0.9

# Right header
add_text(s6, Inches(4.9), Inches(0.3), Inches(8.0), Inches(0.6),
         "解决方案", size=22, bold=True, color=LIGHT_BLUE)

# Right bullets
solutions6 = [
    ("RAG", "通过 RAG 实现知识库，包含培训产品知识，培训话术，考核关键项等知识"),
    ("虚拟人", "实现贴心且恰当的类人对话，提升用户体验"),
    ("LLM", "实现智能陪练智能体，包含意图识别，RAG 知识获取，组织语言答复，考核评价"),
]
y = 1.05
for label, body in solutions6:
    add_text(s6, Inches(4.9), Inches(y), Inches(0.1), Inches(0.3),
             "", fill=LIGHT_BLUE)
    add_rich(s6, Inches(5.05), Inches(y-0.02), Inches(8.0), Inches(0.4), [
        [(label + "：", {"size": 12, "bold": True, "color": WHITE}),
         (body, {"size": 11, "color": WHITE})],
    ])
    y += 0.4

# AI COACH user journey panel
add_text(s6, Inches(4.9), Inches(2.4), Inches(8.2), Inches(4.7),
         "", fill=RGBColor(0x14, 0x1F, 0x4A), line=LIGHT_BLUE)
add_text(s6, Inches(5.05), Inches(2.5), Inches(8.0), Inches(0.4),
         "AI COACH 用户旅程 - 界面", size=14, bold=True, color=WHITE)

# Admin row
add_text(s6, Inches(5.05), Inches(3.0), Inches(0.9), Inches(0.4),
         "Admin", size=10, color=WHITE)
admin_steps = ["数据培训列表", "新建陪练", "设置陪练选项", "配置对话环节", "陪练创建完成"]
sw = 7.2/5
for i, s in enumerate(admin_steps):
    add_text(s6, Inches(6.0 + i*sw), Inches(3.0), Inches(sw-0.05), Inches(0.3),
             s, size=9, color=WHITE, align=PP_ALIGN.CENTER)

# User row
add_text(s6, Inches(5.05), Inches(5.8), Inches(0.9), Inches(0.4),
         "User", size=10, color=WHITE)
user_steps = ["选择培训", "选择素材", "沉浸式 AI 智能陪练", "完成多维评估反馈"]
sw2 = 7.2/4
for i, s in enumerate(user_steps):
    add_text(s6, Inches(6.0 + i*sw2), Inches(5.8), Inches(sw2-0.05), Inches(0.3),
             s, size=9, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s6, Inches(10.0), Inches(6.7), Inches(3.0), Inches(0.3),
         "移动端 AI 练习场景", size=10, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)

# Arrow visual at bottom
arr = s6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                          Inches(5.05), Inches(6.4), Inches(7.9), Inches(0.2))
arr.fill.solid(); arr.fill.fore_color.rgb = LIGHT_BLUE
arr.line.fill.background()

add_text(s6, Inches(10.5), Inches(7.30), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 7: AI 文献搜索
# ============================================================
s7 = prs.slides.add_slide(blank)

# Background split
add_text(s7, Inches(0), Inches(0), Inches(4.5), SLIDE_H, "",
         fill=RGBColor(0xF4, 0xF6, 0xFA))
add_text(s7, Inches(4.5), Inches(0), Inches(8.833), SLIDE_H, "",
         fill=BG_DARK)

# Left header
add_text(s7, Inches(0.4), Inches(0.3), Inches(4.0), Inches(0.7),
         "AI 文献搜索", size=30, bold=True, color=DARK)
add_text(s7, Inches(0.4), Inches(1.05), Inches(4.0), Inches(0.9),
         "从文献检索到全链路赋能销售代表",
         size=14, color=ORANGE)

# Left timeline
timeline7 = [
    ("拜访前", ["拜访计划生成", "拜访资料准备", "AI 智能教练"]),
    ("拜访中", ["业务知识问询解答"]),
    ("拜访后", ["拜访报告智能撰写", "用户肖像测绘"]),
]
y = 2.3
for stage, items in timeline7:
    # circle badge
    c = s7.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(0.4), Inches(y), Inches(0.8), Inches(0.8))
    c.fill.solid(); c.fill.fore_color.rgb = BLUE
    c.line.fill.background()
    tf = c.text_frame; tf.margin_left=tf.margin_right=Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = stage
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "微软雅黑"

    for j, it in enumerate(items):
        add_text(s7, Inches(1.4), Inches(y + j*0.32), Inches(3.0), Inches(0.3),
                 "• " + it, size=11, color=DARK)
    y += 1.5

# Right header
add_text(s7, Inches(4.9), Inches(0.3), Inches(8.0), Inches(0.6),
         "解决方案", size=22, bold=True, color=LIGHT_BLUE)

add_text(s7, Inches(4.9), Inches(1.0), Inches(8.0), Inches(0.4),
         "提示词工程 + RAG", size=14, bold=True, color=WHITE)

bullets7 = [
    "人工智能 + 业务知识库分析推理，理解用户真实意图",
    "无需对文献资料进行人工审核注标签的运维",
    "支持人机互动，查询结果更精准。",
]
y = 1.5
for b in bullets7:
    add_text(s7, Inches(4.9), Inches(y), Inches(8.0), Inches(0.32),
             "• " + b, size=11, color=WHITE)
    y += 0.34

# Process flow 1-5
flow7 = [
    ("1", "自然语言问询"),
    ("2", ""),
    ("3", "交互式提问"),
    ("4", ""),
    ("5", "查看文献并下载"),
]
fw = 8.0/5
add_text(s7, Inches(4.9), Inches(3.0), Inches(2.4), Inches(0.4),
         "给出搜索结果并分析", size=10, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s7, Inches(9.6), Inches(3.0), Inches(2.4), Inches(0.4),
         "追加资料并分析", size=10, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
for i, (num, label) in enumerate(flow7):
    cx = 4.9 + i*fw
    c = s7.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(cx + fw/2 - 0.2), Inches(3.9), Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = LIGHT_BLUE
    c.line.fill.background()
    tf = c.text_frame; tf.margin_left=tf.margin_right=Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "微软雅黑"
    if label:
        add_text(s7, Inches(cx), Inches(4.4), Inches(fw-0.05), Inches(0.4),
                 label, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# AI capability tags
add_text(s7, Inches(4.9), Inches(5.0), Inches(8.0), Inches(0.3),
         "AI 能力支持", size=11, bold=True, color=LIGHT_BLUE)
caps = ["AI 智能体", "提示词工程", "Embedding", "RAG", "语音识别", "图像处理", "内容创作", "数字人"]
cw = 8.0/8
for i, c in enumerate(caps):
    add_text(s7, Inches(4.9 + i*cw + 0.03), Inches(5.35), Inches(cw-0.06), Inches(0.35),
             c, size=9, color=WHITE, fill=LIGHT_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Data support tags
add_text(s7, Inches(4.9), Inches(5.9), Inches(8.0), Inches(0.3),
         "数据支持", size=11, bold=True, color=GREEN)
datas = ["内容平台", "业务知识图谱", "HCP 360", "CRM 数据"]
dw = 8.0/4
for i, c in enumerate(datas):
    add_text(s7, Inches(4.9 + i*dw + 0.03), Inches(6.25), Inches(dw-0.06), Inches(0.35),
             c, size=10, color=WHITE, fill=GREEN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s7, Inches(10.5), Inches(7.30), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 8: AI 业务场景矩阵
# ============================================================
s8 = prs.slides.add_slide(blank)

add_text(s8, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7),
         "AI 业务场景矩阵", size=28, bold=True, color=DARK)

matrix = [
    ("研发 R&D", "研发周期 -40%  |  代码量 +60%", PURPLE, [
        "AI 辅助需求分析（BRD/FSD 自动生成）",
        "架构方案 AI 评估与技术选型",
        "代码生成 / 补全 / Review 自动化",
        "AI 驱动产品设计优化与仿真",
        "智能测试用例生成 & 自动化",
    ]),
    ("生产制造", "质检准确率 99.5%  |  停机 -60%", TEAL, [
        "预测性设备维护（PdM）",
        "视觉质检 AOI — 实时缺陷检测",
        "生产工艺参数 AI 实时优化",
        "数字孪生与柔性生产",
        "能源管理与碳排放 AI 优化",
    ]),
    ("供应与销售", "库存周转 +35%  |  交付准时率 98%", GREEN, [
        "时序 AI 需求预测（滚动预测）",
        "智能补货与库存优化",
        "动态物流路径规划与调度",
        "供应商风险 AI 评估与监控",
        "AI 驱动的定价与促销决策",
    ]),
    ("企业支撑", "人工工时 -70%  |  响应速度 10x", ORANGE, [
        "HR：简历筛选·面试·入职 Agent",
        "财务：智能报销·对账·发票核验",
        "法务：合同审查·合规·风险预警",
        "IT：智能工单·故障自愈·运维",
        "行政：会议室·出行·日程 AI 调度",
    ]),
    ("市场营销", "内容产出 10x  |  转化率 +25%", RGBColor(0xE0, 0x4B, 0x8E), [
        "AI 文案·图片·视频多渠道内容生成",
        "Sora / 可灵 AI 视频 — 产品演示",
        "用户分层与个性化推荐",
        "AI 分析竞品动态与市场洞察",
        "AI 驱动的广告投放优化",
    ]),
    ("客户服务", "7×24 覆盖  |  人工介入 -65%", PURPLE, [
        "多模态 AI 客服（文字·语音·视频）",
        "情感识别与意图理解",
        "FAQ + 知识库智能问答",
        "工单自动分类·优先级·派发",
        "客户满意度 AI 预测与挽回",
    ]),
]

cols = 3
cell_w = 12.5/3
cell_h = 3.0
for i, (title, metric, col, items) in enumerate(matrix):
    r, c = divmod(i, cols)
    x = 0.4 + c*cell_w
    y = 1.0 + r*(cell_h + 0.2)
    # Top color bar with title
    add_text(s8, Inches(x), Inches(y), Inches(cell_w-0.1), Inches(0.45),
             title, size=14, bold=True, color=WHITE, fill=col,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Metric bar
    add_text(s8, Inches(x), Inches(y+0.45), Inches(cell_w-0.1), Inches(0.32),
             metric, size=10, bold=True, color=col, fill=RGBColor(0xF6, 0xF6, 0xFB),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body bullets
    body_text = "\n".join("· " + it for it in items)
    add_text(s8, Inches(x), Inches(y+0.78), Inches(cell_w-0.1), Inches(cell_h-0.8),
             body_text, size=10, color=DARK,
             fill=WHITE, line=RGBColor(0xDD, 0xDD, 0xE5))

add_text(s8, Inches(10.5), Inches(7.30), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 9: AI 整车测试平台
# ============================================================
s9 = prs.slides.add_slide(blank)

# Split background
add_text(s9, Inches(0), Inches(0), Inches(4.8), SLIDE_H, "",
         fill=RGBColor(0xF4, 0xF6, 0xFA))
add_text(s9, Inches(4.8), Inches(0), Inches(8.533), SLIDE_H, "",
         fill=BG_DARK)

# Left title
add_text(s9, Inches(0.4), Inches(0.3), Inches(4.2), Inches(0.7),
         "AI 整车测试平台", size=26, bold=True, color=DARK)
add_text(s9, Inches(0.4), Inches(1.05), Inches(4.2), Inches(0.9),
         "升级研发体系，用 AI 激活汽车研发潜能",
         size=14, color=ORANGE)

# Business value heading
add_text(s9, Inches(0.4), Inches(2.1), Inches(4.2), Inches(0.4),
         "业务价值", size=16, bold=True, color=BLUE)

add_rich(s9, Inches(0.4), Inches(2.55), Inches(4.2), Inches(3.4), [
    [("场景：", {"size": 10, "bold": True, "color": DARK}),
     ("该平台将无缝接入客户现有研发问题追踪系统，捕捉分析车辆全生命周期的技术挑战、市场反馈与投诉。最大限度地减少问题重复出现：利用过去的见解来识别和解决重复出现的挑战，确保新车型的更高可靠性。降低质量风险：通过应用从以往项目中吸取的经验教训，积极主动地减少潜在的质量问题和召回率。",
      {"size": 9, "color": DARK})],
])

add_rich(s9, Inches(0.4), Inches(5.5), Inches(4.2), Inches(1.8), [
    [("效果：", {"size": 10, "bold": True, "color": DARK}),
     ("项目着眼长远，系统成熟后将进化为预测分析平台，提供风险评估与警报。还将助力客户构建企业知识生态，促进协作与数据决策，最终加速开发周期。",
      {"size": 9, "color": DARK})],
])

# Right header
add_text(s9, Inches(5.1), Inches(0.3), Inches(8.0), Inches(0.6),
         "解决方案", size=22, bold=True, color=LIGHT_BLUE)

# UI mockup placeholder
add_text(s9, Inches(5.1), Inches(1.0), Inches(8.0), Inches(2.6),
         "[ Mercedes × Bebest  AI Assistant Chat 界面预览 ]\n\nHi! What do you want to search today?",
         size=11, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         fill=RGBColor(0x14, 0x1F, 0x4A), line=LIGHT_BLUE)

# Architecture diagram
arch_y = 3.8
add_text(s9, Inches(5.1), Inches(arch_y), Inches(8.0), Inches(3.3),
         "", fill=RGBColor(0x14, 0x1F, 0x4A), line=LIGHT_BLUE)

# Dialog module
add_text(s9, Inches(5.25), Inches(arch_y+0.1), Inches(2.0), Inches(0.3),
         "Dialog module", size=9, bold=True, color=WHITE)
add_text(s9, Inches(5.25), Inches(arch_y+0.4), Inches(2.0), Inches(0.35),
         "LLM", size=10, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(5.25), Inches(arch_y+0.8), Inches(2.0), Inches(0.3),
         "Chat guidance", size=9, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(5.25), Inches(arch_y+1.1), Inches(2.0), Inches(0.3),
         "Prompt Engineering", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(5.25), Inches(arch_y+1.4), Inches(1.2), Inches(0.3),
         "Guidance Mgmt", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(6.5), Inches(arch_y+1.4), Inches(0.75), Inches(0.3),
         "Guidance", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(5.25), Inches(arch_y+1.8), Inches(2.0), Inches(0.3),
         "Multi-rounds dialog tree", size=8, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# NLP model
add_text(s9, Inches(7.35), Inches(arch_y+0.4), Inches(0.5), Inches(1.7),
         "NLP model", size=8, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Retrieval module
add_text(s9, Inches(7.95), Inches(arch_y+0.1), Inches(2.7), Inches(0.3),
         "Retrieval module", size=9, bold=True, color=WHITE)
add_text(s9, Inches(7.95), Inches(arch_y+0.4), Inches(2.7), Inches(0.3),
         "Hybrid Retrieval", size=9, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(7.95), Inches(arch_y+0.75), Inches(1.3), Inches(0.3),
         "Preprocessor", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(9.3), Inches(arch_y+0.75), Inches(1.35), Inches(0.3),
         "Text Chunker / Speller", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(7.95), Inches(arch_y+1.1), Inches(2.7), Inches(0.3),
         "Embedding Models", size=9, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Report module
add_text(s9, Inches(7.95), Inches(arch_y+1.55), Inches(2.7), Inches(0.3),
         "Report module", size=9, bold=True, color=WHITE)
add_text(s9, Inches(7.95), Inches(arch_y+1.85), Inches(0.85), Inches(0.3),
         "LLM", size=9, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(8.85), Inches(arch_y+1.85), Inches(1.8), Inches(0.3),
         "Report Prompt", size=9, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(7.95), Inches(arch_y+2.2), Inches(1.3), Inches(0.3),
         "Templates", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(9.3), Inches(arch_y+2.2), Inches(1.35), Inches(0.3),
         "Agents", size=8, color=WHITE, fill=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Data Center
add_text(s9, Inches(10.85), Inches(arch_y+0.1), Inches(2.1), Inches(0.3),
         "Data Center", size=9, bold=True, color=WHITE)
add_text(s9, Inches(10.85), Inches(arch_y+0.4), Inches(2.1), Inches(0.3),
         "Source DB data", size=8, color=DARK, fill=TEAL,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(10.85), Inches(arch_y+0.8), Inches(2.1), Inches(0.3),
         "Preprocess & Convert", size=8, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(10.85), Inches(arch_y+1.15), Inches(2.1), Inches(0.3),
         "Vector store", size=9, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(10.85), Inches(arch_y+1.55), Inches(2.1), Inches(0.3),
         "Configuration", size=9, bold=True, color=WHITE)
add_text(s9, Inches(10.85), Inches(arch_y+1.85), Inches(2.1), Inches(0.3),
         "Business guidance", size=8, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, Inches(10.85), Inches(arch_y+2.25), Inches(2.1), Inches(0.25),
         "Result Store", size=9, bold=True, color=WHITE)
add_text(s9, Inches(10.85), Inches(arch_y+2.55), Inches(2.1), Inches(0.3),
         "Reports", size=8, color=WHITE, fill=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s9, Inches(10.5), Inches(7.30), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 10: AI 知识库
# ============================================================
s10 = prs.slides.add_slide(blank)

# Split background
add_text(s10, Inches(0), Inches(0), Inches(4.5), SLIDE_H, "",
         fill=RGBColor(0xF4, 0xF6, 0xFA))
add_text(s10, Inches(4.5), Inches(0), Inches(8.833), SLIDE_H, "",
         fill=BG_DARK)

# Left header
add_text(s10, Inches(0.4), Inches(0.3), Inches(4.0), Inches(0.7),
         "AI 知识库", size=30, bold=True, color=DARK)
add_text(s10, Inches(0.4), Inches(1.05), Inches(4.0), Inches(1.2),
         "打破企业信息孤岛，驱动知识管理与智能服务双重效能升级",
         size=13, color=ORANGE)

add_text(s10, Inches(0.4), Inches(2.6), Inches(4.0), Inches(0.4),
         "业务价值", size=16, bold=True, color=ORANGE)

bv10 = [
    ("24 小时智能助手", "企业配了个“最强大脑”，员工随时提问（如产品参数、操作流程），秒获准确答案，减少 70% 重复咨询"),
    ("打造“智能知识中枢”", "整合 OA、ERP 等多系统知识，助力企业建立知识大脑，提升服务效率、赋能员工专业化成长"),
]
y = 3.1
for title, body in bv10:
    add_text(s10, Inches(0.4), Inches(y), Inches(0.1), Inches(0.5),
             "", fill=LIGHT_BLUE)
    add_rich(s10, Inches(0.55), Inches(y-0.02), Inches(3.85), Inches(2.0), [
        [(title + "：", {"size": 12, "bold": True, "color": DARK})],
        [(body, {"size": 10, "color": DARK})],
    ])
    y += 1.85

# Right header
add_text(s10, Inches(4.9), Inches(0.3), Inches(8.0), Inches(0.6),
         "解决方案", size=22, bold=True, color=LIGHT_BLUE)

add_rich(s10, Inches(4.9), Inches(1.0), Inches(8.2), Inches(0.4), [
    [("知识库平台：", {"size": 12, "bold": True, "color": WHITE}),
     ("通过构建知识库平台，企业可以通过 AI 大模型技术帮助实现",
      {"size": 11, "color": WHITE})],
])

# Knowledge base UI placeholder
add_text(s10, Inches(4.9), Inches(1.6), Inches(5.5), Inches(2.6),
         "[ 知识库管理界面截图 ]\n\n定义文件上传支持、业务系统动态刷新、打通数据增量",
         size=10, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         fill=RGBColor(0x14, 0x1F, 0x4A), line=LIGHT_BLUE)

add_text(s10, Inches(4.9), Inches(4.3), Inches(5.5), Inches(2.6),
         "[ 对话信息界面截图 ]\n\n沉淀对话信息最佳输入，用于事后单独计费，提升 AI 服务质量",
         size=10, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         fill=RGBColor(0x14, 0x1F, 0x4A), line=LIGHT_BLUE)

# Right side: application scenarios
add_text(s10, Inches(10.55), Inches(1.6), Inches(2.55), Inches(0.4),
         "AI 知识库应用场景", size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
scens = [
    "企业内部产品知识",
    "通过 AI 触摸知识内容",
    "部门内的产品手册",
]
y = 2.2
for s in scens:
    add_text(s10, Inches(10.55), Inches(y), Inches(2.55), Inches(0.55),
             s, size=10, color=WHITE, fill=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.75

# Bottom narrative
add_rich(s10, Inches(4.9), Inches(7.0), Inches(8.2), Inches(0.4), [
    [("场景：", {"size": 9, "bold": True, "color": WHITE}),
     ("通过 AI 把企业内部的知识，转换成 AI 可以学习的语言，将分散的经验转化为可复用的数字资产的核心引擎，能永久锁住人才智慧，缩短 70% 问题解决时间，让企业在人员流动与市场变化中持续降本增效",
      {"size": 8, "color": WHITE})],
])

add_text(s10, Inches(10.5), Inches(7.30), Inches(2.5), Inches(0.2),
         "Company Confidential", size=8, color=GREY, align=PP_ALIGN.RIGHT)


out = "/home/user/Eidy/ppt_output/AI_Cases_and_Offering.pptx"
prs.save(out)
print(f"saved: {out}")
