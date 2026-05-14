# -*- coding: utf-8 -*-
"""星巴克中国 ERP 本地化拆分项目 — 5/18 拜访汇报 PPT 生成器"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 色板 ----------
NAVY   = RGBColor(0x16, 0x2A, 0x4A)   # 主深蓝
TEAL   = RGBColor(0x00, 0x70, 0x7A)   # 凯捷青
GREEN  = RGBColor(0x1E, 0x6B, 0x4F)   # 星巴克绿
GOLD   = RGBColor(0xC8, 0x8A, 0x2E)   # 强调金
RED    = RGBColor(0xB0, 0x3A, 0x2E)   # 风险红
GREY   = RGBColor(0x5B, 0x63, 0x6E)   # 正文灰
LGREY  = RGBColor(0xEC, 0xEE, 0xF1)   # 浅底
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x22, 0x28, 0x33)   # 近黑

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size, color, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {})
        rpr.append(ea)
    ea.set('typeface', font)


def bg(slide, color=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=2, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (txt,size,color,bold)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold)
    return tb


def header(slide, idx, title, sub=None):
    bg(slide)
    # 顶部色条
    box(slide, 0, 0, SW, Inches(0.13), fill=TEAL)
    # 序号块
    box(slide, Inches(0.6), Inches(0.45), Inches(0.62), Inches(0.62), fill=NAVY)
    text(slide, Inches(0.6), Inches(0.45), Inches(0.62), Inches(0.62),
         [[(idx, 24, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Inches(1.4), Inches(0.42), Inches(10.5), Inches(0.5),
         [[(title, 26, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, Inches(1.42), Inches(0.92), Inches(11), Inches(0.35),
             [[(sub, 13, GREY, False)]])
    box(slide, Inches(0.6), Inches(1.3), Inches(12.13), Pt(2), fill=LGREY)
    # 页脚
    text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
         [[("星巴克中国 ERP 本地化拆分项目  |  凯捷  |  5月18日拜访汇报", 9, RGBColor(0xA8,0xAE,0xB6), False)]])


def bullet(slide, x, y, w, items, gap=0.46, size=13, dot=GOLD, lead=None):
    cy = y
    for it in items:
        box(slide, x, cy + Inches(0.07), Inches(0.09), Inches(0.09), fill=dot, shape=MSO_SHAPE.OVAL)
        if isinstance(it, tuple):  # (head, body)
            text(slide, x + Inches(0.22), cy, w - Inches(0.22), Inches(gap),
                 [[(it[0] + "  ", size, INK, True), (it[1], size, GREY, False)]],
                 line_spacing=1.05)
        else:
            text(slide, x + Inches(0.22), cy, w - Inches(0.22), Inches(gap),
                 [[(it, size, GREY, False)]], line_spacing=1.05)
        cy += Inches(gap)
    return cy


# ============================================================
# Slide 1 — 封面
# ============================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(5.0), SW, Inches(0.08), fill=TEAL)
box(s, Inches(0.9), Inches(1.6), Inches(0.16), Inches(1.7), fill=GOLD)
text(s, Inches(1.3), Inches(1.5), Inches(11), Inches(1.0),
     [[("星巴克中国 ERP 本地化拆分项目", 40, WHITE, True)]])
text(s, Inches(1.3), Inches(2.5), Inches(11), Inches(0.7),
     [[("不是换一套软件,是在高速行驶中换一台能独立行军的引擎", 20, RGBColor(0xBFD,0xE0&0xFF,0xE0) if False else RGBColor(0xC9,0xD6,0xE3), False)]])
text(s, Inches(1.3), Inches(3.2), Inches(11), Inches(0.5),
     [[("5月18日 · 上海 · 首轮拜访汇报", 16, GOLD, True)]])
text(s, Inches(1.3), Inches(5.4), Inches(11), Inches(1.2),
     [[("汉得们比的是「配软件」 —— 凯捷比的是「在万店餐饮零售里,真正做过这台引擎拆分的人」", 15, RGBColor(0x9F,0xB0,0xC2), False)]])
text(s, Inches(1.3), Inches(6.6), Inches(11), Inches(0.4),
     [[("凯捷 Capgemini  |  内部汇报材料", 11, RGBColor(0x6E,0x80,0x95), False)]])

# ============================================================
# Slide 2 — 故事线
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "00", "今天怎么讲 —— 一条故事线,五幕", "不谈商务,只回答一个问题:为什么这件事该交给凯捷")
# 主张框
box(s, Inches(0.6), Inches(1.55), Inches(12.13), Inches(1.15), fill=NAVY)
box(s, Inches(0.6), Inches(1.55), Inches(0.13), Inches(1.15), fill=GOLD)
text(s, Inches(0.95), Inches(1.62), Inches(11.6), Inches(1.0),
     [[("核心主张   ", 13, GOLD, True),
       ("星巴克要的不是「换软件」,是博裕入主后,给中国区换上一台能独立行军的引擎。", 16, WHITE, True)],
      [("竞争维度从「谁会配 SAP / Oracle」,切换到「谁真正做过万店餐饮零售的 ERP 分拆」。", 14, RGBColor(0xC9,0xD6,0xE3), False)]],
     line_spacing=1.15)
# 五幕
acts = [
    ("第一幕", "读懂处境", "我们带着对你架构的判断来,不是来听需求的", TEAL),
    ("第二幕", "点破真难点", "这是一个「分拆」项目,不是「实施」项目", NAVY),
    ("第三幕", "我们做过", "麦当劳 / 伊利 —— 同一道题的答卷", GREEN),
    ("第四幕", "怎么控风险", "分阶段解耦 + 系统解耦路线图 + 过渡层", GOLD),
    ("第五幕", "下一步", "一个具体、低成本、2 周可交付的承诺", RED),
]
cw = Inches(2.32); gap = Inches(0.13); x0 = Inches(0.6); y0 = Inches(3.05); ch = Inches(2.7)
for i, (a, t, d, c) in enumerate(acts):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, ch, fill=LGREY)
    box(s, x, y0, cw, Inches(0.62), fill=c)
    text(s, x, y0, cw, Inches(0.62), [[(a, 14, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.12), y0 + Inches(0.78), cw - Inches(0.24), Inches(0.5),
         [[(t, 17, INK, True)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.15), y0 + Inches(1.35), cw - Inches(0.3), Inches(1.2),
         [[(d, 11.5, GREY, False)]], align=PP_ALIGN.CENTER, line_spacing=1.2)
    if i < 4:
        text(s, x + cw - Inches(0.02), y0 + Inches(1.0), gap + Inches(0.1), Inches(0.5),
             [[("→", 16, GOLD, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(6.0), Inches(12.13), Inches(0.5),
     [[("汇报节奏:开场定调 2′ → 一/二幕 10′ → 三幕(麦当劳案例,本场主角)8′ → 四幕 8′ → 五幕 3′ → Q&A", 12, GREY, False)]])

# ============================================================
# Slide 3 — 第一幕 读懂处境
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "01", "第一幕 · 读懂处境 —— 这件事到底难在哪",
       "用星巴克自己的事实证明:我们做过功课")
text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
     [[("难点不是选型,是「高速换引擎」 —— 2026/4/1 交割,业务一天不能停。", 15, INK, True)]])
# 三个结构性约束
box(s, Inches(0.6), Inches(2.0), Inches(7.4), Inches(3.55), fill=LGREY)
text(s, Inches(0.85), Inches(2.15), Inches(7), Inches(0.4),
     [[("三个结构性约束", 16, NAVY, True)]])
bullet(s, Inches(0.9), Inches(2.7), Inches(6.9), [
    ("共用一个 Instance", "EBS 与北美共用同一 Instance —— 拆的是一对「连体」"),
    ("美方不给代码", "仅提供 EBS 端口技术文档,不给克隆、不给源码 —— 不能 lift-and-shift,要重建"),
    ("接口靠文档反推", "11 Inbound / 22 Outbound 逻辑只能靠 RFI 文档还原 —— 文档质量 = 项目最大单点风险"),
    ("数据合规缺口", "现有系统仅覆盖至 2013 年,审计要求保留期需与财务核实口径"),
], gap=0.68, size=12.5)
# 范围边界
box(s, Inches(8.2), Inches(2.0), Inches(4.53), Inches(3.55), fill=WHITE, line=GREEN, line_w=1.5)
text(s, Inches(8.45), Inches(2.15), Inches(4), Inches(0.4),
     [[("我们也读懂了「边界」", 16, GREEN, True)]])
text(s, Inches(8.45), Inches(2.62), Inches(4), Inches(0.4),
     [[("本次切割范围内", 12.5, INK, True)]])
text(s, Inches(8.45), Inches(2.95), Inches(4), Inches(0.4),
     [[("ERP(EBS 核心)· Hyperion · Anaplan · 主数据 HUB", 11.5, GREY, False)]], line_spacing=1.2)
text(s, Inches(8.45), Inches(3.6), Inches(4), Inches(0.4),
     [[("明确不碰(范围外)", 12.5, RED, True)]])
text(s, Inches(8.45), Inches(3.93), Inches(4), Inches(0.9),
     [[("租赁系统 · BOH(2028 另行更换)· Rainbow 销售预测算法(Control Tower 维护)", 11.5, GREY, False)]], line_spacing=1.2)
text(s, Inches(8.45), Inches(4.85), Inches(4), Inches(0.6),
     [[("主动讲清边界 = 懂分寸、不贪大", 11.5, GREEN, True)]])
# 金句
box(s, Inches(0.6), Inches(5.8), Inches(12.13), Inches(0.7), fill=NAVY)
text(s, Inches(0.6), Inches(5.8), Inches(12.13), Inches(0.7),
     [[("「我们不是来听需求的,是带着对你架构的判断来的。」", 16, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 4 — 第二幕 真难点
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "02", "第二幕 · 点破真难点 —— 这是「分拆」项目,不是「实施」项目",
       "把竞争维度从「谁会配软件」换到「谁懂 Carve-out」")
# EBS 三位一体
text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
     [[("Oracle EBS 在星巴克中国是「三位一体」 —— 拆 ERP = 拆整张蜘蛛网", 15, INK, True)]])
roles = [
    ("财务核心", "GL / AP / AR / FA / PA / CE\n法定账套与核算逻辑", TEAL),
    ("集成枢纽", "11 Inbound + 22 Outbound\n前后端系统全靠它串联", NAVY),
    ("主数据中心", "Oracle HUB 分发\nItem / Supplier / Store / COA", GREEN),
]
cw = Inches(3.9); gap = Inches(0.2); x0 = Inches(0.6); y0 = Inches(2.1)
for i, (t, d, c) in enumerate(roles):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, Inches(1.9), fill=WHITE, line=c, line_w=2)
    box(s, x, y0, cw, Inches(0.55), fill=c)
    text(s, x, y0, cw, Inches(0.55), [[(t, 16, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.15), y0 + Inches(0.72), cw - Inches(0.3), Inches(1.1),
         [[(line, 12.5, GREY, False)] for line in d.split("\n")],
         align=PP_ALIGN.CENTER, line_spacing=1.25)
box(s, Inches(0.6), Inches(4.3), Inches(12.13), Inches(1.0), fill=LGREY)
text(s, Inches(0.95), Inches(4.45), Inches(11.5), Inches(0.7),
     [[("真正的能力门槛是「解耦」", 15, NAVY, True),
       ("  —— 不是会不会配 SAP / Oracle,而是能不能把「连体」干净地拆开、还让业务不停。", 13.5, GREY, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
box(s, Inches(0.6), Inches(5.55), Inches(12.13), Inches(0.85), fill=NAVY)
text(s, Inches(0.6), Inches(5.55), Inches(12.13), Inches(0.85),
     [[("「会配软件的供应商很多;做过这种规模分拆的,没几家。」", 16, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 5 — 第三幕 我们做过
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "03", "第三幕 · 我们做过 —— 麦当劳 / 伊利的同题答卷",
       "讲法:不是讲历史,是讲「星巴克现在的风险 = 当年的题,我们这样解的」")
# 对照表
rows = [
    ("星巴克现在的风险", "凯捷做过的同题", TEAL),
    ("EBS 与北美共用一个 Instance", "麦当劳中国从 Global 分拆 —— 共用 Instance 的拆分实操"),
    ("美方只给文档、不给代码", "总部只给接口文档不给源码,凯捷完成本地重建"),
    ("前端门店业务切换期不能停", "POS / BOH 在后端 ERP 切换期间「无感」过渡的过渡层设计"),
    ("万店规模、海量交易", "伊利等大型企业本地架构落地;凯捷中国区 1000+ 人财务服务团队"),
    ("财务报表 / 预算割裂", "财务中台、合并报表、全面预算、票税、银企直连全链路能力"),
]
x0 = Inches(0.6); y0 = Inches(1.55); w1 = Inches(5.0); w2 = Inches(7.13); rh = Inches(0.78)
for i, row in enumerate(rows):
    y = y0 + i * rh
    if i == 0:
        box(s, x0, y, w1, rh, fill=RED)
        box(s, x0 + w1, y, w2, rh, fill=GREEN)
        text(s, x0, y, w1, rh, [[(row[0], 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x0 + w1, y, w2, rh, [[(row[1], 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        fill = WHITE if i % 2 else LGREY
        box(s, x0, y, w1, rh, fill=fill, line=RGBColor(0xDD,0xE0,0xE4), line_w=0.75)
        box(s, x0 + w1, y, w2, rh, fill=fill, line=RGBColor(0xDD,0xE0,0xE4), line_w=0.75)
        text(s, x0 + Inches(0.18), y, w1 - Inches(0.3), rh, [[(row[0], 12.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        text(s, x0 + w1 + Inches(0.18), y, w2 - Inches(0.36), rh, [[(row[1], 12.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
     [[("本页由具备麦当劳分拆背景的顾问主讲 —— 全场可信度的支点。", 12.5, GOLD, True)]])

# ============================================================
# Slide 6 — 第四幕 凯捷打法
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "04", "第四幕 · 怎么控风险 —— 凯捷的打法",
       "给方法论,但落到具体、可追问")
plays = [
    ("1  分阶段解耦", "不一次性切换。并行 → 过渡 → 脱钩,正好对上「两年内继续用 Global」的节奏。", TEAL),
    ("2  系统解耦路线图", "区分「物理拆分」与「逻辑迁移」—— 汉得给不出的硬通货。详见下一节。", NAVY),
    ("3  过渡层设计", "接口与主数据缓冲层,前端门店业务无感,业务连续性优先。", GREEN),
    ("4  财务侧做减法", "Hyperion + Anaplan 走统一平台、统一界面,呼应「消除割裂、避免厂商锁定」。", GOLD),
    ("5  跨文化协调", "凯捷 Global 网络推动美方数据移交谈判,降低本地团队直接交涉风险。", RED),
]
y = Inches(1.6)
for t, d, c in plays:
    box(s, Inches(0.6), y, Inches(0.13), Inches(0.78), fill=c)
    box(s, Inches(0.73), y, Inches(12.0), Inches(0.78), fill=LGREY)
    text(s, Inches(1.0), y, Inches(2.7), Inches(0.78), [[(t, 15, c, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.7), y, Inches(8.8), Inches(0.78), [[(d, 12.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    y += Inches(0.92)
box(s, Inches(0.6), y + Inches(0.05), Inches(12.13), Inches(0.78), fill=NAVY)
text(s, Inches(0.6), y + Inches(0.05), Inches(12.13), Inches(0.78),
     [[("选型不表态:Oracle / SAP 双技术栈都有 —— ", 13.5, RGBColor(0xC9,0xD6,0xE3), False),
       ("「我们帮你把选型决策做对,而不是替某个厂商站台。」", 14, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 7 — 解耦路线图 · 总览
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "04A", "系统解耦路线图 · 总览 —— 三阶段、两条处理路径",
       "本次拜访唯一一张「硬通货」图")
# 三阶段时间轴
phases = [
    ("阶段一 · 解耦", "交割前 ~ 2026/4/1", "建中国区独立 Instance,与 Global 并行;主数据 HUB 本地化;搭过渡层", TEAL),
    ("阶段二 · 过渡", "2026/4/1 ~ +24 个月", "双轨运行,Global 仍在用;接口逐批从「指向北美」切到「指向本地」;数据移交", GOLD),
    ("阶段三 · 脱钩", "~2028", "完全脱离 Global;BOH 等剩余系统按计划更换", GREEN),
]
cw = Inches(3.95); gap = Inches(0.14); x0 = Inches(0.6); y0 = Inches(1.55)
for i, (t, sub, d, c) in enumerate(phases):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, Inches(1.85), fill=WHITE, line=c, line_w=2)
    box(s, x, y0, cw, Inches(0.5), fill=c)
    text(s, x, y0, cw, Inches(0.5), [[(t, 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.12), y0 + Inches(0.58), cw - Inches(0.24), Inches(0.32),
         [[(sub, 11.5, c, True)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.15), y0 + Inches(0.92), cw - Inches(0.3), Inches(0.85),
         [[(d, 11, GREY, False)]], align=PP_ALIGN.CENTER, line_spacing=1.15)
    if i < 2:
        text(s, x + cw - Inches(0.04), y0 + Inches(0.6), gap + Inches(0.12), Inches(0.5),
             [[("→", 18, GOLD, True)]], align=PP_ALIGN.CENTER)
# 两条路径
text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.4),
     [[("每个系统只走两条路径之一 —— 先分类,再排期", 15, INK, True)]])
box(s, Inches(0.6), Inches(4.2), Inches(6.0), Inches(2.05), fill=WHITE, line=RED, line_w=2)
box(s, Inches(0.6), Inches(4.2), Inches(6.0), Inches(0.5), fill=RED)
text(s, Inches(0.6), Inches(4.2), Inches(6.0), Inches(0.5),
     [[("路径 A · 物理拆分", 15, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.85), Inches(4.85), Inches(5.5), Inches(1.3),
     [[("必须真拆 —— 重建中国区独立实例", 12.5, INK, True)],
      [("与北美共用 Instance / 承载法定账套或主数据源头 / 切割后失去美方支持的系统。", 11.5, GREY, False)],
      [("成本高、周期长,是项目的关键路径。", 11.5, RED, True)]], line_spacing=1.2)
box(s, Inches(6.73), Inches(4.2), Inches(6.0), Inches(2.05), fill=WHITE, line=TEAL, line_w=2)
box(s, Inches(6.73), Inches(4.2), Inches(6.0), Inches(0.5), fill=TEAL)
text(s, Inches(6.73), Inches(4.2), Inches(6.0), Inches(0.5),
     [[("路径 B · 逻辑迁移", 15, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(6.98), Inches(4.85), Inches(5.5), Inches(1.3),
     [[("只换归属与连接 —— 业务逻辑不动", 12.5, INK, True)],
      [("通过接口重指向即可解决的周边系统,本体不需重建。", 11.5, GREY, False)],
      [("成本低、可分批,用过渡层兜底。", 11.5, TEAL, True)]], line_spacing=1.2)
text(s, Inches(0.6), Inches(6.45), Inches(12), Inches(0.4),
     [[("收敛判断:33 个接口不是 33 个项目 —— 分类后真正高风险的物理拆分项,可能不到 1/3。", 12.5, GOLD, True)]])

# ============================================================
# Slide 8 — 解耦路线图 · 判断标准
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "04B", "系统解耦路线图 · 判断标准 —— 「物理拆 or 逻辑迁」的四问",
       "任意一问为「是」→ 物理拆分;四问全为「否」→ 逻辑迁移")
qs = [
    ("Q1  连体性", "是否与北美共用同一 Instance / 共享数据库?", "是 → 物理拆分:无法靠接口切断"),
    ("Q2  账套源头", "是否承载法定财务账套,或是主数据的源头系统?", "是 → 物理拆分:必须落到中国区独立实例"),
    ("Q3  支持中断", "切割后美方是否不再提供运维 / 升级支持?", "是 → 物理拆分:不能依赖一个失去支持的系统"),
    ("Q4  接口可解", "是否无法仅靠「接口重指向」解决归属问题?", "是 → 物理拆分;否 → 逻辑迁移即可"),
]
y = Inches(1.6)
for t, q, a in qs:
    box(s, Inches(0.6), y, Inches(2.2), Inches(1.05), fill=NAVY)
    text(s, Inches(0.6), y, Inches(2.2), Inches(1.05), [[(t, 15, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(2.85), y, Inches(5.4), Inches(1.05), fill=LGREY)
    text(s, Inches(3.1), y, Inches(5.0), Inches(1.05), [[(q, 13, INK, True)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    box(s, Inches(8.3), y, Inches(4.43), Inches(1.05), fill=WHITE, line=RED, line_w=1.25)
    text(s, Inches(8.55), y, Inches(4.0), Inches(1.05), [[(a, 12, GREY, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y += Inches(1.2)
box(s, Inches(0.6), y + Inches(0.02), Inches(12.13), Inches(0.62), fill=GREEN)
text(s, Inches(0.6), y + Inches(0.02), Inches(12.13), Inches(0.62),
     [[("这套标准让分类「可解释、可追问」—— 拜访现场被问到任何一个系统,都能当场给出归类逻辑。", 13, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 9 — 解耦路线图 · 系统分类清单
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "04C", "系统解耦路线图 · 系统分类清单(初判,待 RFI 文档校准)",
       "基于现有架构图的初步归类 —— 拿到 RFI 后出精确版")
# 三列
cols = [
    ("路径 A · 物理拆分", RED, [
        "Oracle EBS 财务核心(GL/AP/AR/FA/PA/CE)",
        "Oracle HUB(主数据分发:Item/Supplier/Store/COA)",
        "Hyperion(报表合并,建本地独立 Instance)",
        "Anaplan(全面预算,本地化承接)",
        "Oracle MFG / OPM / QA(烘焙制造)",
    ]),
    ("路径 B · 逻辑迁移", TEAL, [
        "S2P(供应商准入 / 合同 / 订单 / 付款)",
        "3PL WMS(收货 GR / 拣货 Pick Release)",
        "NGBOH(门店进销存)· Control Tower",
        "NGT&E(报销)· SUN(电商)· SIMS(门店管理)",
        "BIP · CDP · TMS · BPM / NGBPM · Vertex",
        "Blackline(对账)· Citi Bank 银企接口",
    ]),
    ("范围外 · 本次不动", GREY, [
        "租赁系统(本次不纳入切割)",
        "BOH 微软系统(2028 年另行更换)",
        "Rainbow 销售预测算法(Control Tower 维护)",
        "Oracle Global 系统(交割后两年内继续使用)",
    ]),
]
cw = Inches(3.95); gap = Inches(0.14); x0 = Inches(0.6); y0 = Inches(1.55); ch = Inches(4.5)
for i, (t, c, items) in enumerate(cols):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, ch, fill=WHITE, line=c, line_w=1.75)
    box(s, x, y0, cw, Inches(0.55), fill=c)
    text(s, x, y0, cw, Inches(0.55), [[(t, 14, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cy = y0 + Inches(0.72)
    for it in items:
        box(s, x + Inches(0.18), cy + Inches(0.06), Inches(0.08), Inches(0.08), fill=c, shape=MSO_SHAPE.OVAL)
        text(s, x + Inches(0.36), cy, cw - Inches(0.5), Inches(0.7),
             [[(it, 11, GREY, False)]], line_spacing=1.1)
        cy += Inches(0.72)
box(s, Inches(0.6), Inches(6.25), Inches(12.13), Inches(0.6), fill=LGREY)
text(s, Inches(0.85), Inches(6.25), Inches(11.6), Inches(0.6),
     [[("注:", 12, RED, True),
       ("本表为基于架构图的初判,部分系统归类需 RFI 文档与现场确认 —— 这正是下一步「接口可拆解性评估」要交付的。", 12, GREY, False)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 10 — 过渡层设计
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "04D", "过渡层设计 —— 让前端门店业务「无感」",
       "业务连续性优先:在高速行驶中换引擎的关键缓冲")
# 简易示意图
# 左:前端
box(s, Inches(0.6), Inches(2.0), Inches(2.6), Inches(3.4), fill=LGREY)
text(s, Inches(0.6), Inches(2.1), Inches(2.6), Inches(0.4), [[("前端 / 门店侧", 13, NAVY, True)]], align=PP_ALIGN.CENTER)
for j, nm in enumerate(["POS", "NGBOH", "3PL WMS", "天猫 / 电商", "S2P"]):
    box(s, Inches(0.8), Inches(2.55) + j*Inches(0.55), Inches(2.2), Inches(0.42), fill=WHITE, line=NAVY, line_w=1)
    text(s, Inches(0.8), Inches(2.55) + j*Inches(0.55), Inches(2.2), Inches(0.42), [[(nm, 11.5, INK, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 中:过渡层
box(s, Inches(4.0), Inches(2.0), Inches(5.0), Inches(3.4), fill=GOLD)
text(s, Inches(4.0), Inches(2.15), Inches(5.0), Inches(0.45), [[("过渡层(缓冲 / 转换)", 15, WHITE, True)]], align=PP_ALIGN.CENTER)
for j, (h, d) in enumerate([
    ("接口缓冲", "前端按原协议收发,过渡层做格式与路由转换"),
    ("主数据双写", "新旧实例同步,HUB 切换期数据一致"),
    ("路由开关", "接口逐批「指向北美 → 指向本地」,可灰度可回滚"),
]):
    box(s, Inches(4.25), Inches(2.65) + j*Inches(0.85), Inches(4.5), Inches(0.72), fill=WHITE)
    text(s, Inches(4.4), Inches(2.65) + j*Inches(0.85), Inches(4.2), Inches(0.72),
         [[(h + "  ", 12, INK, True), (d, 10.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
# 右:后端
box(s, Inches(9.8), Inches(2.0), Inches(2.93), Inches(3.4), fill=LGREY)
text(s, Inches(9.8), Inches(2.1), Inches(2.93), Inches(0.4), [[("后端 ERP 切换", 13, NAVY, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(10.0), Inches(2.6), Inches(2.5), Inches(1.0), fill=WHITE, line=RED, line_w=1.5)
text(s, Inches(10.0), Inches(2.6), Inches(2.5), Inches(1.0), [[("旧:共用 Instance\n(Oracle EBS)", 11, GREY, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
box(s, Inches(10.0), Inches(3.85), Inches(2.5), Inches(1.0), fill=WHITE, line=GREEN, line_w=1.5)
text(s, Inches(10.0), Inches(3.85), Inches(2.5), Inches(1.0), [[("新:中国区独立\nInstance", 11, GREEN, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
# 箭头
for ay in [Inches(3.5)]:
    text(s, Inches(3.2), ay, Inches(0.8), Inches(0.5), [[("→", 22, GOLD, True)]], align=PP_ALIGN.CENTER)
    text(s, Inches(9.0), ay, Inches(0.8), Inches(0.5), [[("→", 22, GOLD, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(0.85), fill=NAVY)
text(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(0.85),
     [[("价值:后端引擎可以「分批换」,前端门店一天都不用停 —— 这是麦当劳分拆验证过的打法。", 14, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 11 — 第五幕 下一步
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "05", "第五幕 · 下一步 —— 一个具体、低成本的承诺",
       "让会议有下文:不空谈,给一个 2 周可交付的钩子")
box(s, Inches(0.6), Inches(1.6), Inches(12.13), Inches(1.7), fill=GREEN)
text(s, Inches(0.95), Inches(1.78), Inches(11.5), Inches(0.5),
     [[("承诺动作", 14, RGBColor(0xCDE&0xFF|0xC0,0xE5,0xD8), True)]])
text(s, Inches(0.95), Inches(2.15), Inches(11.5), Inches(1.0),
     [[("星巴克提供 RFI 文档后,凯捷 2 周内交付", 18, WHITE, True)],
      [("《接口可拆解性评估 + 风险分级清单》—— 把 33 个接口按「物理拆 / 逻辑迁」分类,标出真正高风险的关键路径项。", 13.5, RGBColor(0xDD,0xEE,0xE6), False)]],
     line_spacing=1.25)
# 后续
text(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4), [[("后续配合", 15, NAVY, True)]])
bullet(s, Inches(0.9), Inches(4.05), Inches(11.5), [
    ("案例材料包", "整理 Oracle / SAP 餐饮零售实施实例 + 从 Global Instance 分拆的成功项目复盘"),
    ("专家资源对接", "技术、数据、财务三维度专家按需参与后续会议"),
    ("商务节奏", "今天不谈商务;采购流程按星巴克节奏推进,凯捷不催"),
], gap=0.6, size=13)
box(s, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.7), fill=LGREY)
text(s, Inches(0.95), Inches(6.05), Inches(11.6), Inches(0.7),
     [[("为什么这个承诺有效:", 12.5, GOLD, True),
       ("低成本、可验证、直击星巴克最大单点风险(接口靠文档反推)—— 一次就建立专业信任。", 12.5, GREY, False)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 12 — 收尾金句
# ============================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, Inches(1.0), Inches(2.4), Inches(0.16), Inches(2.5), fill=GOLD)
text(s, Inches(1.4), Inches(2.3), Inches(10.8), Inches(2.8),
     [[("「汉得强在『配软件』,凯捷强在『懂分拆』。", 26, WHITE, True)],
      [("", 10, WHITE, False)],
      [("万店规模的餐饮零售 ERP 切换,我在麦当劳验证过一条铁律 ——", 22, RGBColor(0xC9,0xD6,0xE3), False)],
      [("稳,永远比快值钱。", 26, GOLD, True)],
      [("", 10, WHITE, False)],
      [("我们要交付给星巴克的,不是一套 SAP 或 Oracle,", 22, RGBColor(0xC9,0xD6,0xE3), False)],
      [("是博裕入主后,星巴克中国能独立行军的技术底座。」", 26, WHITE, True)]],
     line_spacing=1.15)
box(s, Inches(1.4), Inches(6.1), Inches(4.0), Pt(2.5), fill=TEAL)
text(s, Inches(1.4), Inches(6.25), Inches(10), Inches(0.5),
     [[("凯捷 Capgemini  ·  星巴克中国 ERP 本地化拆分项目", 13, RGBColor(0x9F,0xB0,0xC2), False)]])

# ============================================================
# Slide 13 — 附:节奏与分工
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "附", "汇报节奏与现场分工", "35 分钟首轮轻量面谈 · 听众:IT 负责人 / PS 团队 / 财务关键用户 / 上级主管")
rows = [
    ("时段", "内容", "主讲", NAVY),
    ("0 – 2′", "开场定调:抛出核心主张", "负责人"),
    ("2 – 12′", "第一、二幕 —— 读懂处境 + 点破真难点", "资深顾问"),
    ("12 – 20′", "第三幕 —— 麦当劳 / 伊利同题答卷", "麦当劳背景顾问(本场主角)"),
    ("20 – 28′", "第四幕 —— 凯捷打法 + 系统解耦路线图", "资深顾问 + 麦当劳背景顾问"),
    ("28 – 31′", "第五幕 —— 下一步承诺", "负责人"),
    ("31 – 35′", "Q&A", "全员"),
]
x0 = Inches(0.6); y0 = Inches(1.6); w = [Inches(2.0), Inches(7.0), Inches(3.13)]; rh = Inches(0.72)
for i, row in enumerate(rows):
    y = y0 + i * rh
    cx = x0
    for j in range(3):
        if i == 0:
            box(s, cx, y, w[j], rh, fill=NAVY)
            text(s, cx, y, w[j], rh, [[(row[j], 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            fill = WHITE if i % 2 else LGREY
            box(s, cx, y, w[j], rh, fill=fill, line=RGBColor(0xDD,0xE0,0xE4), line_w=0.75)
            al = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            bold = (j == 0)
            col = TEAL if j == 0 else (INK if j == 2 else GREY)
            text(s, cx + (Emu(0) if j==0 else Inches(0.2)), y, w[j] - (Emu(0) if j==0 else Inches(0.3)), rh,
                 [[(row[j], 12, col, bold)]], align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += w[j]
box(s, Inches(0.6), Inches(6.0), Inches(12.13), Inches(0.85), fill=LGREY)
text(s, Inches(0.95), Inches(6.0), Inches(11.6), Inches(0.85),
     [[("提醒:", 12.5, RED, True),
       ("采购部门本轮不参与;首轮聚焦背景与案例、建立信任,不进商务细节。第三幕务必由麦当劳背景顾问主讲。", 12.5, GREY, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

prs.save("星巴克ERP拆分_5月18日汇报.pptx")
print("saved:", "星巴克ERP拆分_5月18日汇报.pptx", "slides:", len(prs.slides._sldIdLst))
